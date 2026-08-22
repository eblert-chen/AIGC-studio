# -*- coding: utf-8 -*-
"""生成《零基础读懂项目登录系统》PPT（基于 docs/auth-beginner-guide.md）。"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn

# ---------- 调色板（呼应项目旭日赭橙 + 深色标题栏）----------
INK      = RGBColor(0x1F, 0x29, 0x37)   # 深墨标题栏
SLATE    = RGBColor(0x33, 0x3D, 0x4D)
TEXT     = RGBColor(0x23, 0x2B, 0x36)
MUTED    = RGBColor(0x5B, 0x66, 0x73)
AMBER    = RGBColor(0xB8, 0x55, 0x1D)   # 旭日赭橙 accent
AMBER_LT = RGBColor(0xF6, 0xEC, 0xE4)
PAPER    = RGBColor(0xFF, 0xFF, 0xFF)
PANEL    = RGBColor(0xF4, 0xF6, 0xF8)
LINE     = RGBColor(0xE2, 0xE6, 0xEB)
GREEN    = RGBColor(0x2E, 0x7D, 0x53)
RED      = RGBColor(0xC0, 0x39, 0x2B)

CJ_FONT = "Microsoft YaHei"

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def cjk(run, name=CJ_FONT):
    run.font.name = name
    rPr = run._r.get_or_add_rPr()
    rPr.set(qn('a:ea'), name)


def add_slide():
    return prs.slides.add_slide(BLANK)


def rect(slide, x, y, w, h, fill, line=None, line_w=None):
    from pptx.enum.shapes import MSO_SHAPE
    sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = line_w or Pt(1)
    sp.shadow.inherit = False
    return sp


def textbox(slide, x, y, w, h, lines, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    """lines: list of (text, size, bold, color, space_after)"""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Pt(2); tf.margin_right = Pt(2)
    tf.margin_top = Pt(2); tf.margin_bottom = Pt(2)
    for i, ln in enumerate(lines):
        text, size, bold, color = ln[0], ln[1], ln[2], ln[3]
        sa = ln[4] if len(ln) > 4 else Pt(6)
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = sa
        p.space_before = Pt(0)
        r = p.add_run(); r.text = text
        r.font.size = Pt(size); r.font.bold = bold
        r.font.color.rgb = color
        cjk(r)
    return tb


def header(slide, kicker, title):
    rect(slide, 0, 0, SW, Inches(1.15), INK)
    rect(slide, 0, Inches(1.15), SW, Pt(3), AMBER)
    textbox(slide, Inches(0.55), Inches(0.16), Inches(12), Inches(0.32),
            [(kicker, 12, True, AMBER, 0)])
    textbox(slide, Inches(0.55), Inches(0.44), Inches(12.2), Inches(0.62),
            [(title, 25, True, PAPER, 0)])


def footer(slide, n):
    textbox(slide, Inches(0.55), Inches(7.08), Inches(9), Inches(0.3),
            [("零基础读懂项目登录系统  ·  旭天 AI VIDEO", 9, False, MUTED, 0)])
    textbox(slide, Inches(11.8), Inches(7.08), Inches(1.0), Inches(0.3),
            [(str(n), 9, False, MUTED, 0)], align=PP_ALIGN.RIGHT)


# ============ 1. 封面 ============
s = add_slide()
rect(s, 0, 0, SW, SH, INK)
rect(s, 0, Inches(4.55), SW, Pt(4), AMBER)
textbox(s, Inches(0.9), Inches(1.5), Inches(11.5), Inches(0.5),
        [("AI 视频平台 · 安全架构科普", 15, True, AMBER, 0)])
textbox(s, Inches(0.9), Inches(2.1), Inches(11.5), Inches(1.6),
        [("零基础读懂", 50, True, PAPER, 4),
         ("项目登录系统", 50, True, PAPER, 0)])
textbox(s, Inches(0.9), Inches(4.8), Inches(11.5), Inches(1.4),
        [("OIDC + JWT + Session 三件套，到底各自管什么？", 20, False, RGBColor(0xCF,0xD6,0xDE), 10),
         ("每个名词都配大白话解释 · 附酒店类比", 14, False, RGBColor(0x9A,0xA4,0xB0), 0)])

# ============ 2. 登录解决什么问题 ============
s = add_slide(); header(s, "先说人话", "登录，到底在解决什么问题？")
textbox(s, Inches(0.55), Inches(1.5), Inches(12.2), Inches(0.7),
        [("一个系统每天要回答三个问题——登录系统专门负责前两个。", 16, False, TEXT, 0)])
cards = [
    ("① 你是谁？", "确认身份，确保不是别人冒充你。", AMBER),
    ("② 你已进来吗？", "记住登录状态，别让你每点一下都重输密码。", SLATE),
    ("③ 你能做吗？", "权限检查（这是另一套系统管的事）。", MUTED),
]
cx = Inches(0.55); cw = Inches(3.95); gap = Inches(0.2); cy = Inches(2.5); ch = Inches(2.6)
for i, (t, d, col) in enumerate(cards):
    x = cx + i * (cw + gap)
    rect(s, x, cy, cw, ch, PANEL, line=LINE, line_w=Pt(1))
    rect(s, x, cy, cw, Inches(0.7), col)
    textbox(s, x+Inches(0.2), cy+Inches(0.12), cw-Inches(0.4), Inches(0.5),
            [(t, 18, True, PAPER, 0)])
    textbox(s, x+Inches(0.25), cy+Inches(0.95), cw-Inches(0.5), Inches(1.5),
            [(d, 15, False, TEXT, 0)])
textbox(s, Inches(0.55), Inches(5.5), Inches(12.2), Inches(1.2),
        [("记住一句话：认证 = 你是谁；授权 = 你能干啥。本项目登录只管前两件，",
          15, False, MUTED, 4),
         ("“你能干啥”由权限系统另外管。", 15, False, MUTED, 0)])
footer(s, 2)

# ============ 3. 名词词典（上）============
s = add_slide(); header(s, "名词小词典（上）", "先认识这几个“主角”")
rows = [
    ("认证 Authentication", "证明“你真的是你” —— 酒店前台核身份证"),
    ("授权 Authorization", "确认“你能做某件事” —— 房卡只能开你自己的门"),
    ("身份提供商 IdP", "专门负责“证明你是谁”的外部机构 —— 公安局"),
    ("OIDC", "在 OAuth2 上叠了一层“你是谁”的协议 —— 标准办入住流程"),
    ("OAuth 2.0", "“授权”标准，负责把登录安全地跑完 —— 核验+发卡步骤"),
    ("JWT", "一段被签名、防篡改的文字证明 —— 带防伪的身份证"),
    ("ID Token", "IdP 用 JWT 开的“这是某某人”证明 —— 防伪身份证明"),
    ("JWKS / RS256", "IdP 公开的公钥清单 / 签名算法 —— 验伪设备+印章"),
]
y = Inches(1.45); rh = Inches(0.62)
for i, (k, v) in enumerate(rows):
    yy = y + i * rh
    if i % 2 == 0:
        rect(s, Inches(0.55), yy, Inches(12.2), rh, PANEL)
    textbox(s, Inches(0.7), yy+Inches(0.08), Inches(3.6), rh,
            [(k, 14, True, AMBER, 0)], anchor=MSO_ANCHOR.MIDDLE)
    textbox(s, Inches(4.35), yy+Inches(0.08), Inches(8.2), rh,
            [(v, 13, False, TEXT, 0)], anchor=MSO_ANCHOR.MIDDLE)
footer(s, 3)

# ============ 4. 名词词典（下）============
s = add_slide(); header(s, "名词小词典（下）", "再认识这几个“配角”")
rows = [
    ("Session 会话", "登录后服务器记住“你已登录”的记录 —— 房卡+后台登记"),
    ("Cookie", "浏览器自动保存、访问时自动带上的小数据 —— 随身房卡"),
    ("HttpOnly / Secure", "脚本读不到 / 只在加密连接发 —— 贴身口袋+安全通道"),
    ("SameSite / __Host-", "限同站使用 / 仅主域可设 —— 只在本酒店刷"),
    ("CSRF / CSRF Token", "跨站冒用攻击 / 写操作暗号 —— 开门还要对暗语"),
    ("PKCE / nonce", "登录谜面+答案 / 一次性水印 —— 防偷码+防重用"),
    ("WebAuthn / passkey", "指纹人脸或物理密钥登录 —— 刷脸开门"),
    ("amr / step-up / 吊销", "登录方式记录 / 危险操作再验 / 凭证立失效 —— 登记+金库+挂失"),
]
y = Inches(1.45); rh = Inches(0.62)
for i, (k, v) in enumerate(rows):
    yy = y + i * rh
    if i % 2 == 0:
        rect(s, Inches(0.55), yy, Inches(12.2), rh, PANEL)
    textbox(s, Inches(0.7), yy+Inches(0.08), Inches(3.8), rh,
            [(k, 14, True, AMBER, 0)], anchor=MSO_ANCHOR.MIDDLE)
    textbox(s, Inches(4.55), yy+Inches(0.08), Inches(8.0), rh,
            [(v, 13, False, TEXT, 0)], anchor=MSO_ANCHOR.MIDDLE)
footer(s, 4)

# ============ 5. 核心类比：住酒店 ============
s = add_slide(); header(s, "核心类比", "把登录想成：住酒店")
mapping = [
    ("OIDC / OAuth2", "前台核验身份的流程", "让你出示身份证、走完入住手续", AMBER),
    ("JWT（ID Token）", "身份证本身", "公安局开出的、带防伪签名的“你是谁”证明", SLATE),
    ("Platform Session", "酒店发的房卡 + 后台登记", "确认你住店后给你卡，之后开门刷这张卡", GREEN),
]
cy = Inches(1.6); ch = Inches(1.35); cw = Inches(12.2); gap = Inches(0.18)
for i, (t, sub, d, col) in enumerate(mapping):
    y = cy + i * (ch + gap)
    rect(s, Inches(0.55), y, cw, ch, PANEL, line=LINE, line_w=Pt(1))
    rect(s, Inches(0.55), y, Inches(0.18), ch, col)
    textbox(s, Inches(0.95), y+Inches(0.14), Inches(3.5), ch,
            [(t, 18, True, col, 2), (sub, 13, False, MUTED, 0)],
            anchor=MSO_ANCHOR.MIDDLE)
    textbox(s, Inches(4.7), y+Inches(0.14), Inches(7.8), ch,
            [(d, 15, False, TEXT, 0)], anchor=MSO_ANCHOR.MIDDLE)
rect(s, Inches(0.55), Inches(6.05), Inches(12.2), Inches(0.7), AMBER_LT)
textbox(s, Inches(0.75), Inches(6.12), Inches(11.8), Inches(0.6),
        [("关键句：你不会每次开门都重新出示身份证；用户也不会每次调用平台都重新登录。",
          15, True, AMBER, 0)], anchor=MSO_ANCHOR.MIDDLE)
footer(s, 5)

# ============ 6. 一次完整登录（9步）============
s = add_slide(); header(s, "完整流程", "一次登录，经历了什么？（按时间）")
steps = [
    "你点“登录” → 浏览器告诉平台：我要登录",
    "平台把你转向 IdP，记下 state（防掉包）和 nonce（防重用）",
    "你在 IdP 登录（输密码或刷脸）—— 平台完全不碰你的密码",
    "IdP 把你转回平台，附带一个“授权码”code",
    "平台拿 code 去换 ID Token（JWT 格式，RS256 签名）",
    "平台验 ID Token：用 JWKS 公钥核对签名、iss、aud、nonce",
    "平台创建 Session：数据库写“已登录”记录（只存摘要，不存原文）",
    "平台把房卡交给浏览器：写入安全 Cookie",
    "之后每次操作，浏览器自动带 Cookie，平台查库确认有效就放行",
]
col1 = steps[:5]; col2 = steps[5:]
for ci, blk in enumerate([col1, col2]):
    x = Inches(0.55) + ci * Inches(6.3)
    for i, st in enumerate(blk):
        yy = Inches(1.55) + i * Inches(1.02)
        rect(s, x, yy, Inches(0.5), Inches(0.5), AMBER)
        textbox(s, x, yy, Inches(0.5), Inches(0.5),
                [(str(i+1), 16, True, PAPER, 0)], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        rect(s, x+Inches(0.62), yy, Inches(5.55), Inches(0.9), PANEL, line=LINE, line_w=Pt(1))
        textbox(s, x+Inches(0.78), yy+Inches(0.06), Inches(5.3), Inches(0.8),
                [(st, 13, False, TEXT, 0)], anchor=MSO_ANCHOR.MIDDLE)
textbox(s, Inches(0.55), Inches(6.75), Inches(12.2), Inches(0.4),
        [("注意第 5–7 步：JWT 验完即弃，浏览器手里只有房卡（Session Cookie），没有长期 JWT。",
          13, True, AMBER, 0)])
footer(s, 6)

# ============ 7. 为什么用三种？============
s = add_slide(); header(s, "为什么不是一种？", "只用其中一种，会怎样？")
cols = [
    ("只用 JWT", RED, [
        "签发后在过期前一直有效",
        "“退出全部设备”旧 JWT 可能还能用",
        "要立即吊销就得建黑名单 → 又变回 Session",
        "放 localStorage 易被 XSS 偷走",
    ]),
    ("只用 Session", SLATE, [
        "平台得自建完整身份系统",
        "自己管密码加密/找回/邮箱",
        "自己管 MFA/WebAuthn/企业 SSO",
        "密钥轮换、异常防护…安全责任巨大",
    ]),
    ("只用 OIDC", AMBER, [
        "浏览器要保存 IdP Token",
        "与 IdP 生命周期强耦合",
        "难独立管理设备/权限变更",
        "平台全局停用，IdP Token 可能仍有效",
    ]),
]
cw = Inches(3.95); gap = Inches(0.2); cx = Inches(0.55); cy = Inches(1.5); ch = Inches(4.7)
for i, (t, col, items) in enumerate(cols):
    x = cx + i * (cw + gap)
    rect(s, x, cy, cw, ch, PANEL, line=LINE, line_w=Pt(1))
    rect(s, x, cy, cw, Inches(0.6), col)
    textbox(s, x, cy+Inches(0.1), cw, Inches(0.45),
            [(t, 17, True, PAPER, 0)], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    lines = [(("•  " + it), 13, False, TEXT, Pt(10)) for it in items]
    textbox(s, x+Inches(0.25), cy+Inches(0.8), cw-Inches(0.5), ch-Inches(1.0), lines)
footer(s, 7)

# ============ 8. 三件套组合的好处 ============
s = add_slide(); header(s, "组合的好处", "三件套一起，同时拿到这些")
benefits = [
    "企业统一登录（SSO）", "WebAuthn / passkey 强认证", "JWT / JWKS 签名验证（防伪）",
    "浏览器不保存 IdP Token（XSS 偷不到）", "当前设备立即退出", "全部设备立即退出",
    "平台全局停用立即生效", "独立的空闲 + 绝对有效期", "高风险操作要求重新认证 (step-up)",
    "前端无法靠改用户 ID 冒充别人",
]
cy = Inches(1.55); cw = Inches(3.95); ch = Inches(0.95); gx = Inches(0.2); gy = Inches(0.2)
for i, b in enumerate(benefits):
    r, c = divmod(i, 3)
    x = Inches(0.55) + c * (cw + gx)
    y = cy + r * (ch + gy)
    rect(s, x, y, cw, ch, AMBER_LT, line=LINE, line_w=Pt(1))
    rect(s, x, y, Inches(0.12), ch, GREEN)
    textbox(s, x+Inches(0.28), y, cw-Inches(0.4), ch,
            [(b, 14, True, TEXT, 0)], anchor=MSO_ANCHOR.MIDDLE)
footer(s, 8)

# ============ 9. 每次请求被查的 8 关 ============
s = add_slide(); header(s, "后台在偷偷干啥", "登录后，每次请求被查 8 关")
checks = [
    "房卡在不在？ —— Cookie 里的 Session Token 有吗",
    "房卡是真的吗？ —— 数据库里有对应摘要记录吗",
    "这张卡挂失了吗？ —— Session 是否被撤销",
    "人还在吗？ —— 用户账号是否仍 ACTIVE",
    "版本对吗？ —— auth_version 匹配（改密/权限后旧卡失效）",
    "过期了吗？ —— 绝对有效期 + 空闲有效期",
    "危险操作要再验？ —— auth_time 够新否则弹窗重登 (step-up)",
    "管理员刷脸了吗？ —— 高风险请求要求 amr 含 WebAuthn/passkey",
]
y = Inches(1.5); rh = Inches(0.62)
for i, c in enumerate(checks):
    yy = y + i * rh
    if i % 2 == 0:
        rect(s, Inches(0.55), yy, Inches(12.2), rh, PANEL)
    textbox(s, Inches(0.7), yy+Inches(0.06), Inches(0.5), rh,
            [(str(i+1), 14, True, AMBER, 0)], anchor=MSO_ANCHOR.MIDDLE)
    textbox(s, Inches(1.25), yy+Inches(0.06), Inches(11.3), rh,
            [(c, 13, False, TEXT, 0)], anchor=MSO_ANCHOR.MIDDLE)
textbox(s, Inches(0.55), Inches(6.55), Inches(12.2), Inches(0.5),
        [("因为每次都查数据库 Session，所以“停用账号 / 退出全部设备”能立刻生效。",
          14, True, GREEN, 0)])
footer(s, 9)

# ============ 10. 房卡的安全属性 ============
s = add_slide(); header(s, "房卡上的防伪标签", "Session Cookie 的安全属性")
attrs = [
    ("Secure", "只在加密连接(https)下发送，明文网络里不出现"),
    ("HttpOnly", "不能用 JavaScript 读取，防 XSS 偷走"),
    ("SameSite=Lax", "别的网站发起的请求不带它，防 CSRF 冒用"),
    ("Path=/", "全站有效"),
    ("__Host- 前缀", "只能由主域名设置，防子域冒名下发"),
    ("CSRF Token", "写操作随请求带暗号，平台比对一致才放行"),
]
y = Inches(1.5); rh = Inches(0.78)
for i, (k, v) in enumerate(attrs):
    yy = y + i * rh
    rect(s, Inches(0.55), yy, Inches(3.2), Inches(0.66), INK)
    textbox(s, Inches(0.7), yy, Inches(3.0), Inches(0.66),
            [(k, 15, True, PAPER, 0)], anchor=MSO_ANCHOR.MIDDLE)
    rect(s, Inches(3.8), yy, Inches(8.95), Inches(0.66), PANEL, line=LINE, line_w=Pt(1))
    textbox(s, Inches(4.0), yy, Inches(8.6), Inches(0.66),
            [(v, 13, False, TEXT, 0)], anchor=MSO_ANCHOR.MIDDLE)
footer(s, 10)

# ============ 11. 代码在哪 ============
s = add_slide(); header(s, "想翻代码？", "在你的项目里，它们对应哪里")
code_rows = [
    ("登录 / 回调路由（OAuth2 授权码 + PKCE 起始）", "routers/authentication.py : 224"),
    ("PKCE / nonce / 登录事务生成", "services/authentication.py : 517"),
    ("服务端 Session 创建与验证（查摘要、查撤销）", "services/authentication.py : 208"),
    ("Session 数据模型（存摘要不存原文）", "models.py : 280"),
]
y = Inches(1.6); rh = Inches(0.95)
for i, (desc, loc) in enumerate(code_rows):
    yy = y + i * rh
    rect(s, Inches(0.55), yy, Inches(8.6), Inches(0.8), PANEL, line=LINE, line_w=Pt(1))
    textbox(s, Inches(0.8), yy, Inches(8.2), Inches(0.8),
            [(desc, 14, False, TEXT, 0)], anchor=MSO_ANCHOR.MIDDLE)
    rect(s, Inches(9.25), yy, Inches(3.5), Inches(0.8), AMBER)
    textbox(s, Inches(9.25), yy, Inches(3.5), Inches(0.8),
            [(loc, 13, True, PAPER, 0)], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
textbox(s, Inches(0.55), Inches(5.9), Inches(12.2), Inches(1.0),
        [("还有一套旧的 HS256 Bearer JWT 登录入口，仅供本地开发，生产环境明确禁止启用——",
          14, False, MUTED, 4),
         ("它不是正式方案，别被它迷惑。", 14, False, MUTED, 0)])
footer(s, 11)

# ============ 12. 一句话总结 ============
s = add_slide()
rect(s, 0, 0, SW, SH, INK)
rect(s, 0, Inches(3.4), SW, Pt(4), AMBER)
textbox(s, Inches(0.9), Inches(1.2), Inches(11.5), Inches(0.5),
        [("一句话总结", 16, True, AMBER, 0)])
textbox(s, Inches(0.9), Inches(1.9), Inches(11.5), Inches(1.4),
        [("OIDC 是前台核身份证，", 30, True, PAPER, 6),
         ("JWT 是那张带防伪的证明，", 30, True, PAPER, 6),
         ("Session 是酒店发给你的房卡。", 30, True, PAPER, 0)])
textbox(s, Inches(0.9), Inches(3.8), Inches(11.5), Inches(2.6),
        [("三者不是重复建设，而是前后衔接：", 18, False, RGBColor(0xCF,0xD6,0xDE), 10),
         ("IdP 只出现在登录边界，验完身份，", 18, False, RGBColor(0xCF,0xD6,0xDE), 6),
         ("平台用自己的房卡（Session）管你之后的每一次访问。", 18, False, RGBColor(0xCF,0xD6,0xDE), 14),
         ("这样你既享受专业身份机构的安全，又拥有平台可控的会话管理，", 16, False, RGBColor(0x9A,0xA4,0xB0), 6),
         ("还把“浏览器偷 Token”的风险压到最低。", 16, False, RGBColor(0x9A,0xA4,0xB0), 0)])

prs.save("docs/auth-beginner.pptx")
print("SAVED docs/auth-beginner.pptx slides=", len(prs.slides._sldIdLst))
